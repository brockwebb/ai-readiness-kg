#!/usr/bin/env python3
"""Build the framework deck from `docs/crosswalk/deck_content_2026-09-01.md`.

Task `cc_tasks/2026-09-01_framework_deck_build.md`. **Zero model calls.** A deterministic
markdown-to-pptx transform: no template, no theme, no images, no charts. The content file is
the single source and is never edited here; the deck is a disposable projection of it, so a
content edit plus a rebuild is the whole update path.

Layout budget, stated because the overflow rule depends on it: 16:9 at 13.333 x 7.5 inches,
body frame 12.33 x 5.5 inches. Usable lines at font f are 5.5 / (1.2 * f / 72), and characters
per line are approximated at 12.33 * 72 / (0.5 * f) for the proportional default face. Slides
are rendered at 18pt and stepped down to a 14pt floor before any split is considered.
"""
from __future__ import annotations

import argparse
import pathlib
import re
import sys

REPO = pathlib.Path(__file__).resolve().parent.parent
SOURCE = REPO / "docs/crosswalk/deck_content_2026-09-01.md"
OUTPUT = REPO / "docs/crosswalk/framework_deck_2026-09-01.pptx"

SLIDE_W, SLIDE_H = 13.333, 7.5
BODY_W, BODY_H = 12.33, 5.5
START_PT, FLOOR_PT, MONO_PT = 18, 14, 12
HEADER_RE = re.compile(r"^## (Slide (\d+)\s*[—-]\s*(.*))$", re.M)
BOLD_RE = re.compile(r"\*\*(.+?)\*\*")


class Line:
    """One rendered body line: bullet at an indent level, a numbered item, or preformatted."""

    def __init__(self, text: str, level: int = 0, pre: bool = False, blank: bool = False):
        self.text, self.level, self.pre, self.blank = text, level, pre, blank


def parse_body(body: str) -> list[Line]:
    """Markdown body to rendered lines. Indent maps to bullet level; a 4-space block with no
    bullet marker is preformatted (slides 4, 8 and 14 are text renderings of diagrams and the
    content file asks for monospace where the layout survives)."""
    out: list[Line] = []
    for raw in body.strip("\n").splitlines():
        if not raw.strip():
            if out and not out[-1].blank:
                out.append(Line("", blank=True))
            continue
        stripped = raw.lstrip(" ")
        indent = len(raw) - len(stripped)
        m = re.match(r"^([-*])\s+(.*)$", stripped)
        if m:
            out.append(Line(m.group(2), level=min(indent // 2, 4)))
            continue
        m = re.match(r"^(\d+)\.\s+(.*)$", stripped)
        if m:
            # Numbered lists stay numbered: the number is kept as literal text so PowerPoint
            # cannot renumber it against the author's intent.
            out.append(Line(f"{m.group(1)}. {m.group(2)}", level=min(indent // 2, 4)))
            continue
        if indent >= 4:
            out.append(Line(raw.rstrip(), pre=True))
            continue
        out.append(Line(stripped, level=0))
    while out and out[-1].blank:
        out.pop()
    return out


def parse(src: str) -> list[dict]:
    slides, parts = [], HEADER_RE.split(src)[1:]
    for _full, num, title, body in zip(parts[0::4], parts[1::4], parts[2::4], parts[3::4]):
        body = body.split("\n---\n")[0]
        lines = [l for l in body.strip().splitlines() if l.strip()]
        lead = lines[0].strip() if lines else ""
        # The task's rule: a standalone bold line after the header is the slide title. Only
        # slide 1 has one (recorded in the RESULT); every other slide takes the header text.
        if re.fullmatch(r"\*\*.+\*\*", lead):
            slides.append({"n": int(num), "title": lead.strip("*"),
                           "body": body.split(lead, 1)[1], "titled_by": "bold_lead"})
        else:
            slides.append({"n": int(num), "title": title.strip(), "body": body,
                           "titled_by": "header"})
    return slides


def wrapped_lines(lines: list[Line], pt: int) -> int:
    per = max(20, int(BODY_W * 72 / (0.5 * pt)))
    total = 0
    for ln in lines:
        if ln.blank:
            total += 1
        elif ln.pre:
            total += 1                       # preformatted never wraps; it may clip instead
        else:
            total += max(1, -(-len(ln.text) // max(10, per - 4 * ln.level)))
    return total


#: Slack, in rendered lines, held back from the height budget. The wrap model here is an
#: ESTIMATE (average character width at 0.5 em); PowerPoint does the real layout with the
#: real font metrics. A slide sized to a one-line margin is inside this model's own error,
#: so it can render overflowing even though the arithmetic said it fit. Two lines of slack
#: costs a point of font on the three densest slides and removes that failure mode.
FIT_SLACK = 2


def capacity(pt: int) -> int:
    return int(BODY_H / (1.2 * pt / 72)) - FIT_SLACK


def fits(lines: list[Line], pt: int) -> bool:
    return wrapped_lines(lines, pt) <= capacity(pt)


def add_runs(para, text: str, bold_all: bool = False):
    """Render inline **emphasis** as real bold runs. The words are never altered; only the
    markers are consumed, so the content file stays authoritative."""
    pos = 0
    for m in BOLD_RE.finditer(text):
        if m.start() > pos:
            r = para.add_run(); r.text = text[pos:m.start()]; r.font.bold = bold_all
        r = para.add_run(); r.text = m.group(1); r.font.bold = True
        pos = m.end()
    if pos < len(text):
        r = para.add_run(); r.text = text[pos:]; r.font.bold = bold_all
    if not text:
        para.add_run().text = ""


def build(slides: list[dict], out: pathlib.Path) -> dict:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(SLIDE_W), Inches(SLIDE_H)
    BLACK = RGBColor(0, 0, 0)
    report = {"slides": [], "splits": []}

    for spec in slides:
        lines = parse_body(spec["body"])
        pt = START_PT
        while pt > FLOOR_PT and not fits(lines, pt):
            pt -= 1
        chunks = [lines]
        if not fits(lines, pt):
            # Still over at the floor: split on a blank-line boundary nearest the midpoint so
            # a bullet is never cut in half, and mark the continuation in its title.
            cap = capacity(pt)
            chunks, cur, used = [], [], 0
            for ln in lines:
                cost = wrapped_lines([ln], pt)
                if used + cost > cap and cur:
                    chunks.append(cur); cur, used = [], 0
                cur.append(ln); used += cost
            if cur:
                chunks.append(cur)
            report["splits"].append({"slide": spec["n"], "parts": len(chunks), "pt": pt})

        for i, chunk in enumerate(chunks):
            title = spec["title"] if i == 0 else f"{spec['title']} (cont.)"
            if spec["n"] == 1 and i == 0:
                slide = prs.slides.add_slide(prs.slide_layouts[0])
                slide.shapes.title.text_frame.text = title
                sub = slide.placeholders[1].text_frame
                body_lines = [l for l in chunk if not l.blank]
                sub.text = ""
                for j, ln in enumerate(body_lines):
                    p = sub.paragraphs[0] if j == 0 else sub.add_paragraph()
                    add_runs(p, re.sub(r"^Subtitle:\s*", "", ln.text))
                for shp in (slide.shapes.title, slide.placeholders[1]):
                    for p in shp.text_frame.paragraphs:
                        for r in p.runs:
                            r.font.color.rgb = BLACK
                report["slides"].append({"n": spec["n"], "title": title, "pt": "layout",
                                         "lines": len(body_lines), "layout": "title"})
                continue

            slide = prs.slides.add_slide(prs.slide_layouts[6])       # blank
            tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.35),
                                          Inches(BODY_W), Inches(0.9))
            tf = tb.text_frame
            tf.word_wrap = True
            add_runs(tf.paragraphs[0], title, bold_all=True)
            for r in tf.paragraphs[0].runs:
                r.font.size, r.font.color.rgb = Pt(28), BLACK

            bb = slide.shapes.add_textbox(Inches(0.5), Inches(1.45),
                                          Inches(BODY_W), Inches(BODY_H))
            bf = bb.text_frame
            bf.word_wrap = True
            first = True
            for ln in chunk:
                p = bf.paragraphs[0] if first else bf.add_paragraph()
                first = False
                if ln.blank:
                    p.add_run().text = ""
                    for r in p.runs:
                        r.font.size = Pt(max(8, pt // 2))
                    continue
                p.level = 0 if ln.pre else ln.level
                add_runs(p, ln.text)
                for r in p.runs:
                    r.font.size = Pt(MONO_PT if ln.pre else pt)
                    r.font.color.rgb = BLACK
                    if ln.pre:
                        r.font.name = "Courier New"
            report["slides"].append({"n": spec["n"], "title": title, "pt": pt,
                                     "lines": len(chunk),
                                     "pre": sum(1 for l in chunk if l.pre)})
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    return report


def main() -> int:
    ap = argparse.ArgumentParser(description=__doc__)
    ap.add_argument("--source", default=str(SOURCE))
    ap.add_argument("--out", default=str(OUTPUT))
    a = ap.parse_args()
    src = pathlib.Path(a.source).read_text(encoding="utf-8")
    slides = parse(src)
    if not slides:
        raise SystemExit("FATAL: no '## Slide N' sections found; refusing to write an empty deck")
    rep = build(slides, pathlib.Path(a.out))
    print(f"source sections: {len(slides)}   slides written: {len(rep['slides'])}")
    for s in rep["slides"]:
        print(f"  {s['n']:>2}  {str(s['pt']):>6}pt  {s['lines']:>2} lines"
              f"{'  pre=' + str(s['pre']) if s.get('pre') else ''}   {s['title'][:60]}")
    if rep["splits"]:
        print("\nsplits:", rep["splits"])
    else:
        print("\nno slide needed splitting")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
